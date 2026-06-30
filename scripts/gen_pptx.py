"""Generate capabilities-deck.pptx from HTML deck content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0f, 0x17, 0x2a)
BG_CARD   = RGBColor(0x1e, 0x29, 0x3b)
BLUE      = RGBColor(0x3b, 0x82, 0xf6)
INDIGO    = RGBColor(0x63, 0x66, 0xf1)
GREEN     = RGBColor(0x22, 0xc5, 0x5e)
RED       = RGBColor(0xef, 0x44, 0x44)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
ADESSO    = RGBColor(0xe3, 0x00, 0x0f)
WHITE     = RGBColor(0xf8, 0xfa, 0xfc)
LIGHT     = RGBColor(0xcb, 0xd5, 0xe1)
MID       = RGBColor(0x94, 0xa3, 0xb8)
DIM       = RGBColor(0x64, 0x74, 0x8b)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide

def txb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def tf_para(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
            space_before=0, space_after=0, italic=False):
    para = tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after  = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return para

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_logo(slide, x=Inches(12.6), y=Inches(7.05), size=Pt(14)):
    tb = txb(slide, x, y, Inches(0.7), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "adesso"
    run.font.bold  = True
    run.font.size  = size
    run.font.color.rgb = ADESSO

def add_slide_number(slide, n, x=Inches(0.3), y=Inches(7.05)):
    tb = txb(slide, x, y, Inches(0.5), Inches(0.35))
    tf = tb.text_frame
    para = tf.add_paragraph()
    run = para.add_run()
    run.text = str(n)
    run.font.size  = Pt(10)
    run.font.color.rgb = DIM

def add_accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, 0, W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

def eyebrow(slide, text, left, top, width=Inches(8)):
    tb = txb(slide, left, top, width, Inches(0.25))
    tf = tb.text_frame
    para = tf.add_paragraph()
    run = para.add_run()
    run.text = text.upper()
    run.font.size  = Pt(9)
    run.font.bold  = True
    run.font.color.rgb = BLUE

def heading(slide, text, left, top, width=Inches(10), size=24):
    tb = txb(slide, left, top, width, Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = WHITE

def bullet_box(slide, items, left, top, width, height,
               dot_color=BLUE, font_size=11, gap=Inches(0.08)):
    tb = txb(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        para = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        para.space_before = Pt(3)
        run = para.add_run()
        run.text = f"› {item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = LIGHT

def card(slide, left, top, width, height, title, body_lines,
         title_color=BLUE, top_accent=None):
    add_rect(slide, left, top, width, height,
             fill_color=BG_CARD, line_color=RGBColor(0x2d,0x3f,0x55))
    if top_accent:
        add_rect(slide, left, top, width, Pt(3), fill_color=top_accent)
    # title
    tb = txb(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.22))
    tf = tb.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = title.upper()
    run.font.size  = Pt(8)
    run.font.bold  = True
    run.font.color.rgb = title_color
    # body
    tb2 = txb(slide, left + Inches(0.15), top + Inches(0.35), width - Inches(0.3), height - Inches(0.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    first = True
    for line in body_lines:
        para2 = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        para2.space_before = Pt(2)
        run2 = para2.add_run()
        run2.text = f"› {line}" if not line.startswith("›") else line
        run2.font.size  = Pt(9)
        run2.font.color.rgb = MID

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)

# adesso logo top-left
tb = txb(s, Inches(0.4), Inches(0.25), Inches(1.6), Inches(0.5))
tf = tb.text_frame
para = tf.paragraphs[0]
run = para.add_run()
run.text = "adesso"
run.font.size  = Pt(28)
run.font.bold  = True
run.font.color.rgb = ADESSO

# RFP ref eyebrow
eyebrow(s, "RFP Response · MC-2026-0417", Inches(0.55), Inches(1.8))

# Main title
tb = txb(s, Inches(0.55), Inches(2.15), Inches(9), Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
para = tf.paragraphs[0]
run = para.add_run()
run.text = "Modernizing Meridian's\nInventory Dashboard"
run.font.size  = Pt(36)
run.font.bold  = True
run.font.color.rgb = WHITE

# Blue divider line
add_rect(s, Inches(0.55), Inches(3.65), Inches(0.5), Pt(3), fill_color=BLUE)

# Subtitle
tb = txb(s, Inches(0.55), Inches(3.85), Inches(8.5), Inches(0.9))
tf = tb.text_frame
tf.word_wrap = True
para = tf.paragraphs[0]
run = para.add_run()
run.text = "adesso is ready to remediate, extend, and future-proof the platform Meridian's previous vendor left behind — on a fixed-fee engagement with no surprises."
run.font.size  = Pt(13)
run.font.color.rgb = MID

# Meta row
meta = [("Prepared by", "adesso"), ("Submitted", "May 8, 2026"), ("Contact", "marcel.cossijns@adesso.de")]
for i, (label, value) in enumerate(meta):
    tb = txb(s, Inches(0.55 + i*3.5), Inches(4.85), Inches(3.4), Inches(0.3))
    tf = tb.text_frame
    para = tf.paragraphs[0]
    r1 = para.add_run(); r1.text = label + " "; r1.font.size = Pt(11); r1.font.color.rgb = DIM
    r2 = para.add_run(); r2.text = value;        r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = MID

add_logo(s)
add_slide_number(s, 1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What Meridian has today
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Current State", Inches(0.5), Inches(0.35))
heading(s, "What Meridian has today", Inches(0.5), Inches(0.65))

issues = [
    "Reports module filters are unwired — data does not filter on user input",
    "No automated test coverage delivered — IT team is blocking all changes",
    "i18n gaps leave Tokyo warehouse staff without translated UI",
    "Inconsistent API patterns across views; console errors in production",
    "No architecture documentation — new engineers cannot onboard",
    "Previous vendor departed November 2024 with minimal handoff",
]
bullet_box(s, issues, Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.8),
           dot_color=RED, font_size=11)

card(s, Inches(6.7), Inches(1.5), Inches(5.9), Inches(1.6),
     "Missing entirely",
     ["Restocking recommendations view", "Budget-ceiling purchase order logic",
      "End-to-end browser tests", "Current-state architecture docs"])

card(s, Inches(6.7), Inches(3.3), Inches(5.9), Inches(1.9),
     "Risk if left unaddressed",
     ["IT gatekeeper blocks deployments. Operations team makes restock decisions without tooling.",
      "Tokyo warehouse operates on a partially translated interface."],
     title_color=RED)

add_logo(s); add_slide_number(s, 2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Stakeholders
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Stakeholder Alignment", Inches(0.5), Inches(0.35))
heading(s, "Who we're working for — and what they need", Inches(0.5), Inches(0.65))

stakeholders = [
    ("VP Operations", "R. Tanaka", AMBER,
     "Daily Dashboard User",
     "Frustrated with the previous vendor. Her priority is the Restocking view (R2) — she needs data-driven purchase recommendations with a budget ceiling her team can adjust."),
    ("IT (gatekeeper)", "Unnamed contact", RED,
     "Blocking Deployments",
     "Has blocked all changes because no test coverage was delivered. R3 (automated browser tests) is non-negotiable — it must land before IT will approve anything."),
    ("Director of Procurement", "J. Okafor", BLUE,
     "Vendor Selection Owner",
     "Owns the decision. Values timeline and price predictability above all. We respond with a fixed-fee structure and a 10-week delivery plan with clear phase gates."),
]

col_w = Inches(4.1)
for i, (role, name, accent, badge, need) in enumerate(stakeholders):
    x = Inches(0.4 + i * 4.3)
    y = Inches(1.55)
    h = Inches(4.5)
    add_rect(s, x, y, col_w, h, fill_color=BG_CARD, line_color=RGBColor(0x2d,0x3f,0x55))

    tb = txb(s, x+Inches(0.15), y+Inches(0.12), col_w-Inches(0.3), Inches(0.22))
    tf = tb.text_frame; para = tf.paragraphs[0]
    run = para.add_run(); run.text = role.upper()
    run.font.size = Pt(8); run.font.bold = True; run.font.color.rgb = DIM

    tb2 = txb(s, x+Inches(0.15), y+Inches(0.38), col_w-Inches(0.3), Inches(0.35))
    tf2 = tb2.text_frame; para2 = tf2.paragraphs[0]
    run2 = para2.add_run(); run2.text = name
    run2.font.size = Pt(14); run2.font.bold = True; run2.font.color.rgb = WHITE

    # badge
    add_rect(s, x+Inches(0.15), y+Inches(0.82), col_w-Inches(0.3), Inches(0.22),
             fill_color=accent)
    tb3 = txb(s, x+Inches(0.18), y+Inches(0.83), col_w-Inches(0.36), Inches(0.2))
    tf3 = tb3.text_frame; para3 = tf3.paragraphs[0]
    run3 = para3.add_run(); run3.text = badge
    run3.font.size = Pt(8); run3.font.bold = True; run3.font.color.rgb = WHITE

    tb4 = txb(s, x+Inches(0.15), y+Inches(1.15), col_w-Inches(0.3), Inches(3.0))
    tf4 = tb4.text_frame; tf4.word_wrap = True; para4 = tf4.paragraphs[0]
    run4 = para4.add_run(); run4.text = need
    run4.font.size = Pt(10); run4.font.color.rgb = MID

tb = txb(s, Inches(0.4), Inches(6.3), Inches(12), Inches(0.3))
tf = tb.text_frame; para = tf.paragraphs[0]
run = para.add_run()
run.text = "3 warehouses · San Francisco (HQ) · London (EMEA) · Tokyo (APAC, opened 2023) · ~180 employees · ~$9.6M revenue"
run.font.size = Pt(10); run.font.color.rgb = DIM

add_logo(s); add_slide_number(s, 3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Approach overview
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Approach Overview", Inches(0.5), Inches(0.35))
heading(s, "Four required deliverables, one structured engagement", Inches(0.5), Inches(0.65))

deliverables = [
    ("R1", "Reports Remediation", BLUE,
     ["Audit every issue, fix at least 8 — filters, i18n gaps, API inconsistencies, console noise."],
     "Phase 1–2"),
    ("R2", "Restocking View", INDIGO,
     ["New Vue view. Stock levels + demand forecast + operator budget ceiling → purchase order recommendations."],
     "Phase 2–3"),
    ("R3", "Automated Tests", GREEN,
     ["Playwright end-to-end coverage for critical flows. CI-ready. Unblocks IT approval gate."],
     "Phase 2–3"),
    ("R4", "Architecture Docs", AMBER,
     ["Current-state overview delivered in Week 2. Enables IT onboarding and safe future changes."],
     "Phase 1"),
]

col_w = Inches(2.9)
for i, (req, title, accent, body, phase) in enumerate(deliverables):
    x = Inches(0.45 + i * 3.1)
    y = Inches(1.55)
    h = Inches(4.4)
    add_rect(s, x, y, col_w, h, fill_color=BG_CARD, line_color=RGBColor(0x2d,0x3f,0x55))
    add_rect(s, x, y, col_w, Pt(3), fill_color=accent)

    tb = txb(s, x+Inches(0.12), y+Inches(0.12), col_w-Inches(0.24), Inches(0.2))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = req; run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = accent

    tb2 = txb(s, x+Inches(0.12), y+Inches(0.35), col_w-Inches(0.24), Inches(0.4))
    tf2 = tb2.text_frame; tf2.word_wrap = True; para2 = tf2.paragraphs[0]; run2 = para2.add_run()
    run2.text = title; run2.font.size = Pt(12); run2.font.bold = True; run2.font.color.rgb = WHITE

    tb3 = txb(s, x+Inches(0.12), y+Inches(0.85), col_w-Inches(0.24), Inches(2.5))
    tf3 = tb3.text_frame; tf3.word_wrap = True; para3 = tf3.paragraphs[0]; run3 = para3.add_run()
    run3.text = body[0]; run3.font.size = Pt(10); run3.font.color.rgb = MID

    tb4 = txb(s, x+Inches(0.12), y+Inches(3.8), col_w-Inches(0.24), Inches(0.3))
    tf4 = tb4.text_frame; para4 = tf4.paragraphs[0]; run4 = para4.add_run()
    run4.text = phase; run4.font.size = Pt(9); run4.font.bold = True; run4.font.color.rgb = accent

tb5 = txb(s, Inches(0.45), Inches(6.3), Inches(12.0), Inches(0.3))
tf5 = tb5.text_frame; para5 = tf5.paragraphs[0]; run5 = para5.add_run()
run5.text = "Plus desired deliverables D1–D3 (UI modernization, full i18n, dark mode) scoped in optional Phase 3 extension."
run5.font.size = Pt(10); run5.font.color.rgb = DIM

add_logo(s); add_slide_number(s, 4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — R1 Reports Remediation
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Required Deliverable · R1", Inches(0.5), Inches(0.35))
heading(s, "Reports Remediation — audit first, fix second", Inches(0.5), Inches(0.65))

issues_r1 = [
    "Filter behavior — dropdowns and date pickers do not wire to the underlying data",
    "i18n gaps — untranslated strings across the Reports view",
    "API pattern inconsistencies — mix of direct fetch and abstracted service calls",
    "Console noise — unhandled promise rejections and Vue warnings in production builds",
]
bullet_box(s, issues_r1, Inches(0.5), Inches(1.55), Inches(5.8), Inches(4.0), font_size=12)

card(s, Inches(6.7), Inches(1.55), Inches(5.9), Inches(2.2),
     "Our method",
     ["Structured audit pass — document each defect before touching code",
      "Deliver audit log to Meridian for sign-off",
      "Fix in order of stakeholder impact",
      "Each fix accompanied by a regression test"])

card(s, Inches(6.7), Inches(3.9), Inches(5.9), Inches(1.7),
     "Commitment",
     ["RFP requires 'at least 8 issues.' We will document everything we find and fix in priority order — no cherry-picking easy wins."],
     title_color=GREEN)

add_logo(s); add_slide_number(s, 5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — R2 Restocking
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Required Deliverable · R2", Inches(0.5), Inches(0.35))
heading(s, "Restocking Recommendations — the feature Tanaka needs", Inches(0.5), Inches(0.65))

items_r2 = [
    "New Vue 3 Composition API view — consistent with existing codebase patterns",
    "Input: current stock levels per SKU + demand forecast signals",
    "Operator-adjustable budget ceiling — the team sets the constraint, the system recommends within it",
    "Output: ranked purchase order recommendations with quantity and estimated cost",
    "FastAPI endpoint backing the calculation — same JSON-file data layer, no DB migration required",
]
bullet_box(s, items_r2, Inches(0.5), Inches(1.55), Inches(5.8), Inches(4.2), font_size=11)

card(s, Inches(6.7), Inches(1.55), Inches(5.9), Inches(1.6),
     "Design constraint",
     ["We work within the existing Vue + FastAPI stack. No new runtime dependencies introduced without client approval. The view will be navigable from the existing sidebar."])

card(s, Inches(6.7), Inches(3.3), Inches(5.9), Inches(2.3),
     "Assumptions (to confirm with Tanaka)",
     ["Budget ceiling is per warehouse, not global",
      "Demand forecast is based on historical order data in the JSON files",
      "Recommendations are advisory, not auto-submitted"])

add_logo(s); add_slide_number(s, 6)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — R3 Automated Testing
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Required Deliverable · R3", Inches(0.5), Inches(0.35))
heading(s, "Automated Browser Testing — unblocking IT", Inches(0.5), Inches(0.65))

flows = [
    "Inventory browse — warehouse filter applied",
    "Orders management — view and status filter",
    "Reports module — filter interactions, export, i18n rendering",
    "Restocking view — budget input, recommendation output, interaction",
]
bullet_box(s, flows, Inches(0.5), Inches(1.55), Inches(5.8), Inches(3.2), dot_color=GREEN, font_size=12)

tb = txb(s, Inches(0.5), Inches(4.9), Inches(5.8), Inches(0.4))
tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
run.text = "Tests are written against localhost:3000 with the full stack running — real network calls, real data files."
run.font.size = Pt(10); run.font.color.rgb = DIM

card(s, Inches(6.7), Inches(1.55), Inches(5.9), Inches(1.7),
     "Technology",
     ["Playwright — the industry standard for browser automation. Test files committed to the repo. CI-ready from day one."],
     title_color=GREEN)

card(s, Inches(6.7), Inches(3.4), Inches(5.9), Inches(2.2),
     "What IT gets",
     ["A passing test suite they can run before approving any PR",
      "Tests as living documentation of expected behavior",
      "A clear path to adding coverage for future features"])

add_logo(s); add_slide_number(s, 7)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — R4 Architecture Documentation
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Required Deliverable · R4", Inches(0.5), Inches(0.35))
heading(s, "Architecture Documentation — delivered in Week 2", Inches(0.5), Inches(0.65))

arch_items = [
    "Current-state overview of the Vue 3 + FastAPI stack — components, routes, data flow",
    "Data layer diagram — how JSON files in server/data/ map to API endpoints and frontend views",
    "Deployment topology — frontend dev server (3000), API server (8001), and how they connect",
    "Known gaps and technical debt noted explicitly — no whitewashing the handoff quality",
    "Delivered as an HTML interactive diagram — shareable without tooling",
]
bullet_box(s, arch_items, Inches(0.5), Inches(1.55), Inches(5.8), Inches(4.0), font_size=11)

card(s, Inches(6.7), Inches(1.55), Inches(5.9), Inches(2.0),
     "Why Week 2?",
     ["Delivering R4 early gives IT and new engineers the map they need to understand every change we make from Week 3 onward. It de-risks the rest of the engagement."],
     title_color=AMBER)

card(s, Inches(6.7), Inches(3.7), Inches(5.9), Inches(1.6),
     "Format",
     ["HTML with embedded SVG diagram. No Visio, no Confluence dependency.",
      "Version-controlled alongside the code. Always up-to-date."])

add_logo(s); add_slide_number(s, 8)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Delivery Plan", Inches(0.5), Inches(0.35))
heading(s, "10 weeks · 3 phases · fixed gates", Inches(0.5), Inches(0.65))

phases = [
    ("Phase 1", "Weeks 1–2",  BLUE,
     ["R4 Architecture Docs", "R1 Audit", "Codebase orientation", "Kickoff & access"]),
    ("Phase 2", "Weeks 3–6",  INDIGO,
     ["R1 Fixes (≥8 issues)", "R2 Restocking View", "R3 Test foundation", "Phase gate review"]),
    ("Phase 3", "Weeks 7–9",  GREEN,
     ["R3 Full test suite", "D1 UI modernization", "D2 i18n (Tokyo)", "D3 Dark mode"]),
    ("Buffer",  "Week 10",    MID,
     ["Review & sign-off", "Contingency", "Final handoff"]),
]

row_h = Inches(1.1)
for i, (name, weeks, color, items) in enumerate(phases):
    y = Inches(1.55 + i * 1.2)
    # label
    tb = txb(s, Inches(0.4), y+Inches(0.1), Inches(1.2), Inches(0.3))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = name; run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = WHITE
    tb2 = txb(s, Inches(0.4), y+Inches(0.45), Inches(1.2), Inches(0.25))
    tf2 = tb2.text_frame; para2 = tf2.paragraphs[0]; run2 = para2.add_run()
    run2.text = weeks; run2.font.size = Pt(9); run2.font.color.rgb = DIM

    # bar
    bar_x = Inches(1.75); bar_w = Inches(11.0)
    add_rect(s, bar_x, y, bar_w, Inches(0.85), fill_color=BG_CARD)
    # items
    item_x = bar_x + Inches(0.15)
    for j, item in enumerate(items):
        tb3 = txb(s, item_x + Inches(j*2.6), y+Inches(0.22), Inches(2.5), Inches(0.4))
        tf3 = tb3.text_frame; para3 = tf3.paragraphs[0]; run3 = para3.add_run()
        run3.text = item; run3.font.size = Pt(9); run3.font.bold = True
        run3.font.color.rgb = LIGHT if color != MID else MID

add_logo(s); add_slide_number(s, 9)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Pricing
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Commercial Terms", Inches(0.5), Inches(0.35))
heading(s, "Fixed-fee · not-to-exceed · no surprises", Inches(0.5), Inches(0.65))

rows = [
    ("R4 Architecture documentation", "2", "€2,500", False),
    ("R1 Reports remediation",         "8", "€10,000", False),
    ("R3 Automated browser testing",   "5", "€6,250",  False),
    ("R2 Restocking recommendations",  "10","€12,500", False),
    ("Project management & delivery",  "5", "€6,250",  False),
    ("Required total (R1–R4)",         "30","€37,500", True),
    ("[Optional extension D1–D3]",     "",  "",         False),
    ("D1 UI modernization",            "5", "€6,250",  False),
    ("D2 Internationalization",        "4", "€5,000",  False),
    ("D3 Dark mode",                   "3", "€3,750",  False),
    ("Full scope NTE",                 "42","€52,500", True),
]

y_start = Inches(1.45)
row_h   = Inches(0.42)
for i, (label, days, fee, is_total) in enumerate(rows):
    y = y_start + i * row_h
    bg = RGBColor(0x1a, 0x2e, 0x4a) if is_total else (
         RGBColor(0x16, 0x23, 0x35) if label.startswith("[") else BG_DARK)
    add_rect(s, Inches(0.4), y, Inches(8.6), row_h, fill_color=bg)

    color = WHITE if is_total else (INDIGO if label.startswith("[") else LIGHT)

    tb = txb(s, Inches(0.5), y+Inches(0.1), Inches(5.5), row_h-Inches(0.05))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = label.replace("[","").replace("]","")
    run.font.size = Pt(10); run.font.bold = is_total; run.font.color.rgb = color

    if days:
        tb2 = txb(s, Inches(6.1), y+Inches(0.1), Inches(1.0), row_h-Inches(0.05))
        tf2 = tb2.text_frame; para2 = tf2.paragraphs[0]; para2.alignment = PP_ALIGN.RIGHT
        run2 = para2.add_run(); run2.text = days
        run2.font.size = Pt(10); run2.font.bold = is_total; run2.font.color.rgb = color

    if fee:
        tb3 = txb(s, Inches(7.2), y+Inches(0.1), Inches(1.7), row_h-Inches(0.05))
        tf3 = tb3.text_frame; para3 = tf3.paragraphs[0]; para3.alignment = PP_ALIGN.RIGHT
        run3 = para3.add_run(); run3.text = fee
        run3.font.size = Pt(10); run3.font.bold = True; run3.font.color.rgb = BLUE if is_total else color

# two footnote cards
card(s, Inches(9.2), Inches(1.55), Inches(3.9), Inches(2.0),
     "Payment milestones",
     ["30% on contract signing",
      "40% at Phase 2 gate",
      "30% on final delivery and sign-off"])

card(s, Inches(9.2), Inches(3.7), Inches(3.9), Inches(2.0),
     "Fixed-fee rationale",
     ["Okafor values predictability. Our fixed fee means Meridian cannot overspend regardless of what we find in the codebase. No change-order creep on required scope."],
     title_color=BLUE)

add_logo(s); add_slide_number(s, 10)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Relevant Experience
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Firm Credentials", Inches(0.5), Inches(0.35))
heading(s, "Relevant Experience", Inches(0.5), Inches(0.65))

projects = [
    ("Project Fluent", "Mid-market industrial distributor, Germany", BLUE,
     ["Vue 3", "FastAPI", "Playwright"],
     "10 weeks · 3 engineers",
     ["Inherited dashboard · defect remediation + test coverage",
      "11 defects resolved across filter, i18n, and formatting",
      "Playwright suite · 8 critical flows · 94% coverage at delivery",
      "Architecture docs handed to IT team"],
     "IT team approved 2 follow-on change requests within Q1. Zero regressions in 6 months."),
    ("Project Apex", "EMEA wholesale distributor, Netherlands", INDIGO,
     ["Vue 3", "FastAPI", "PostgreSQL"],
     "14 weeks · 4 engineers",
     ["Net-new restocking recommendations view on live app",
      "Budget-ceiling purchase order engine (FastAPI)",
      "ERP export workflow (SAP B1 REST integration)",
      "Playwright regression suite for new + existing views"],
     "30% reduction in manual reorder prep time. Delivered on schedule."),
    ("Project Kanji", "Precision parts manufacturer, Germany & Tokyo", GREEN,
     ["Vue 3", "vue-i18n v9", "Playwright"],
     "16 weeks · 4 engineers",
     ["UI modernization + Japanese locale rollout (all modules)",
      "Dark-mode theme for 12 warehouse floor stations",
      "87 Playwright tests across DE/EN/JP locale variants",
      "WCAG 2.1 AA compliance confirmed on primary views"],
     "Tokyo team fully adopted Japanese-locale views within first week."),
]

col_w = Inches(4.1)
for i, (name, client, accent, tags, meta, bullets, outcome) in enumerate(projects):
    x = Inches(0.4 + i * 4.3)
    y = Inches(1.55)
    h = Inches(5.0)
    add_rect(s, x, y, col_w, h, fill_color=BG_CARD, line_color=RGBColor(0x2d,0x3f,0x55))
    add_rect(s, x, y, col_w, Pt(3), fill_color=accent)

    cy = y + Inches(0.12)
    tb = txb(s, x+Inches(0.12), cy, col_w-Inches(0.24), Inches(0.22))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = name; run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = accent

    cy += Inches(0.26)
    tb2 = txb(s, x+Inches(0.12), cy, col_w-Inches(0.24), Inches(0.35))
    tf2 = tb2.text_frame; tf2.word_wrap = True; para2 = tf2.paragraphs[0]; run2 = para2.add_run()
    run2.text = client; run2.font.size = Pt(10); run2.font.bold = True; run2.font.color.rgb = WHITE

    cy += Inches(0.38)
    tb3 = txb(s, x+Inches(0.12), cy, col_w-Inches(0.24), Inches(0.2))
    tf3 = tb3.text_frame; para3 = tf3.paragraphs[0]; run3 = para3.add_run()
    run3.text = "  ".join(tags); run3.font.size = Pt(9); run3.font.bold = True; run3.font.color.rgb = BLUE

    cy += Inches(0.24)
    tb4 = txb(s, x+Inches(0.12), cy, col_w-Inches(0.24), Inches(0.2))
    tf4 = tb4.text_frame; para4 = tf4.paragraphs[0]; run4 = para4.add_run()
    run4.text = meta; run4.font.size = Pt(9); run4.font.color.rgb = DIM

    cy += Inches(0.26)
    tb5 = txb(s, x+Inches(0.12), cy, col_w-Inches(0.24), Inches(2.1))
    tf5 = tb5.text_frame; tf5.word_wrap = True
    first = True
    for b in bullets:
        p = tf5.paragraphs[0] if first else tf5.add_paragraph()
        first = False
        p.space_before = Pt(2)
        r = p.add_run(); r.text = f"› {b}"; r.font.size = Pt(9); r.font.color.rgb = MID

    add_rect(s, x+Inches(0.12), y+h-Inches(0.65), col_w-Inches(0.24), Inches(0.5),
             fill_color=RGBColor(0x0d, 0x25, 0x1a))
    tb6 = txb(s, x+Inches(0.18), y+h-Inches(0.62), col_w-Inches(0.36), Inches(0.45))
    tf6 = tb6.text_frame; tf6.word_wrap = True; para6 = tf6.paragraphs[0]; run6 = para6.add_run()
    run6.text = outcome; run6.font.size = Pt(9); run6.font.bold = True; run6.font.color.rgb = GREEN

tb7 = txb(s, Inches(0.4), Inches(6.7), Inches(12.0), Inches(0.25))
tf7 = tb7.text_frame; para7 = tf7.paragraphs[0]; run7 = para7.add_run()
run7.text = "Client names anonymized per standard reference policy. Detailed case studies and referee contacts available on request."
run7.font.size = Pt(9); run7.font.italic = True; run7.font.color.rgb = DIM

add_logo(s); add_slide_number(s, 11)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Team
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Our Proposed Team", Inches(0.5), Inches(0.35))
heading(s, "Four people. One engagement. No hand-offs.", Inches(0.5), Inches(0.65))

team = [
    ("Lars Meijering", "Engagement Lead",
     ["12 yrs consulting", "PMP certified", "Delivery governance"],
     "Led 3 comparable dashboard modernization engagements. Primary point of contact for J. Okafor and R. Tanaka throughout the engagement."),
    ("Sophie Berger", "Senior Frontend Engineer",
     ["Vue 3 specialist", "7 yrs experience", "vue-i18n contributor"],
     "Open-source contributor to vue-i18n. Built 2 operator-facing inventory UIs in comparable stack. Will lead R1, R2, and D1–D3."),
    ("Daniel Kovač", "Backend & API Engineer",
     ["Python / FastAPI", "6 yrs experience", "API design"],
     "Delivered the restocking algorithm for a wholesale distributor in Project Apex. Owns data pipeline reliability and the R2 backend endpoint."),
    ("Yuki Tanabe", "QA & Test Automation",
     ["Playwright", "Cypress", "CI/CD pipelines", "Native Japanese"],
     "Delivered 94% E2E coverage on Project Fluent. Native Japanese speaker — directly validates Tokyo i18n and locale rendering."),
]

col_w = Inches(6.0)
for i, (name, role, skills, bio) in enumerate(team):
    row = i // 2; col = i % 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.55 + row * 2.5)
    h = Inches(2.2)
    add_rect(s, x, y, col_w, h, fill_color=BG_CARD, line_color=RGBColor(0x2d,0x3f,0x55))

    tb = txb(s, x+Inches(0.15), y+Inches(0.12), col_w-Inches(0.3), Inches(0.3))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = name; run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE

    tb2 = txb(s, x+Inches(0.15), y+Inches(0.46), col_w-Inches(0.3), Inches(0.22))
    tf2 = tb2.text_frame; para2 = tf2.paragraphs[0]; run2 = para2.add_run()
    run2.text = role; run2.font.size = Pt(9); run2.font.bold = True; run2.font.color.rgb = BLUE

    tb3 = txb(s, x+Inches(0.15), y+Inches(0.72), col_w-Inches(0.3), Inches(0.22))
    tf3 = tb3.text_frame; para3 = tf3.paragraphs[0]; run3 = para3.add_run()
    run3.text = "  ·  ".join(skills); run3.font.size = Pt(9); run3.font.color.rgb = MID

    tb4 = txb(s, x+Inches(0.15), y+Inches(1.0), col_w-Inches(0.3), Inches(1.1))
    tf4 = tb4.text_frame; tf4.word_wrap = True; para4 = tf4.paragraphs[0]; run4 = para4.add_run()
    run4.text = bio; run4.font.size = Pt(10); run4.font.color.rgb = DIM

# footer note
add_rect(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.45),
         fill_color=RGBColor(0x10, 0x1f, 0x35))
tb5 = txb(s, Inches(0.55), Inches(6.45), Inches(12.0), Inches(0.35))
tf5 = tb5.text_frame; para5 = tf5.paragraphs[0]; run5 = para5.add_run()
run5.text = "Full CVs available on request. Team allocation: 100% dedicated for Phases 1–2, 50% in Phase 3."
run5.font.size = Pt(10); run5.font.color.rgb = DIM

add_logo(s); add_slide_number(s, 12)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Why adesso
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Firm Credentials", Inches(0.5), Inches(0.35))
heading(s, "Why adesso", Inches(0.5), Inches(0.65))

proof_points = [
    ("Vue 3 + FastAPI depth",
     "We work in this exact stack daily. Composition API, Vite, Python async patterns — no learning curve on the fundamentals."),
    ("Audit-and-fix methodology",
     "We document before we touch. Meridian gets a full defect log before any code changes — not a bill for hours of undocumented work."),
    ("Test-first culture",
     "Playwright is standard practice for us. IT's gating requirement is something we'd have recommended anyway."),
    ("Structured delivery, not heroics",
     "Phase gates, written sign-offs, and a fixed fee. Okafor gets predictability; Tanaka gets her feature; IT gets their tests."),
    ("Multi-region awareness",
     "We understand the Tokyo warehouse context. i18n and timezone handling are first-class concerns, not afterthoughts."),
    ("Evaluation criteria alignment",
     "40% technical approach · 25% experience · 20% timeline · 15% price. This proposal is structured to score well on what Meridian weighted highest."),
]

col_w = Inches(6.0)
for i, (title, body) in enumerate(proof_points):
    row = i % 3; col = i // 3
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.55 + row * 1.6)

    add_rect(s, x+Inches(0.05), y+Inches(0.05), Inches(0.38), Inches(0.38),
             fill_color=RGBColor(0x1a, 0x2e, 0x4a), line_color=BLUE)

    tb = txb(s, x+Inches(0.55), y+Inches(0.05), col_w-Inches(0.65), Inches(0.32))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = title; run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = WHITE

    tb2 = txb(s, x+Inches(0.55), y+Inches(0.4), col_w-Inches(0.65), Inches(1.1))
    tf2 = tb2.text_frame; tf2.word_wrap = True; para2 = tf2.paragraphs[0]; run2 = para2.add_run()
    run2.text = body; run2.font.size = Pt(10); run2.font.color.rgb = DIM

    add_rect(s, x, y+Inches(1.35), col_w, Pt(0.5), fill_color=RGBColor(0x1e, 0x29, 0x3b))

add_logo(s); add_slide_number(s, 13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Next Steps
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_accent_bar(s)
eyebrow(s, "Next Steps", Inches(0.5), Inches(0.35))
heading(s, "How we move from proposal to engagement", Inches(0.5), Inches(0.65))

steps = [
    ("Clarification call with procurement (J. Okafor)",
     "Per RFP §6 — we've submitted 5 written questions. A 30-minute call to confirm budget range, 'critical flows' definition, and Tokyo i18n scope."),
    ("Capabilities presentation (if shortlisted)",
     "We're ready to present within 5 business days of shortlist notification. This deck is our starting point."),
    ("Contract and access provisioning",
     "On award: SOW sign-off, repo access, and a kickoff meeting with Tanaka and IT within one week of contract execution."),
    ("Week 2 architecture doc delivery",
     "First tangible deliverable — gives Meridian something in hand early and signals how we work."),
]

for i, (title, body) in enumerate(steps):
    y = Inches(1.55 + i * 1.15)
    # circle number
    add_rect(s, Inches(0.4), y+Inches(0.05), Inches(0.38), Inches(0.38), fill_color=BLUE)
    tb0 = txb(s, Inches(0.4), y+Inches(0.06), Inches(0.38), Inches(0.3))
    tf0 = tb0.text_frame; para0 = tf0.paragraphs[0]; para0.alignment = PP_ALIGN.CENTER
    run0 = para0.add_run(); run0.text = str(i+1)
    run0.font.size = Pt(11); run0.font.bold = True; run0.font.color.rgb = WHITE

    tb = txb(s, Inches(0.95), y+Inches(0.04), Inches(11.8), Inches(0.3))
    tf = tb.text_frame; para = tf.paragraphs[0]; run = para.add_run()
    run.text = title; run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = WHITE

    tb2 = txb(s, Inches(0.95), y+Inches(0.38), Inches(11.8), Inches(0.6))
    tf2 = tb2.text_frame; tf2.word_wrap = True; para2 = tf2.paragraphs[0]; run2 = para2.add_run()
    run2.text = body; run2.font.size = Pt(10); run2.font.color.rgb = DIM

    add_rect(s, Inches(0.4), y+Inches(1.02), Inches(12.4), Pt(0.5),
             fill_color=RGBColor(0x1e, 0x29, 0x3b))

# contact box
add_rect(s, Inches(0.4), Inches(6.1), Inches(12.4), Inches(0.55),
         fill_color=RGBColor(0x10, 0x1f, 0x35), line_color=RGBColor(0x2d,0x5a,0x9e))
tb3 = txb(s, Inches(0.55), Inches(6.18), Inches(12.0), Inches(0.4))
tf3 = tb3.text_frame; para3 = tf3.paragraphs[0]
r1 = para3.add_run(); r1.text = "Questions or to schedule the clarification call: "
r1.font.size = Pt(11); r1.font.color.rgb = MID
r2 = para3.add_run(); r2.text = "marcel.cossijns@adesso.de"
r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = BLUE
r3 = para3.add_run(); r3.text = "  ·  Reference RFP MC-2026-0417"
r3.font.size = Pt(11); r3.font.color.rgb = MID

add_logo(s); add_slide_number(s, 14)

# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
out = "/home/cossijns/aiw/meridian-workshop/proposal/capabilities-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
